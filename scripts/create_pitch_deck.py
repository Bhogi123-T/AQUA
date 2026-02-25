from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_college_demo_deck():
    prs = Presentation()

    def set_slide_background(slide, color_rgb):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_rgb

    # Colors
    AQUA_BLUE = RGBColor(0, 105, 148)
    DEEP_SEA = RGBColor(0, 25, 50)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BLUE = RGBColor(224, 247, 250)
    ACCENT = RGBColor(0, 255, 136)

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, DEEP_SEA)
    
    title = slide.shapes.title
    title.text = "Aquasphere"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True

    subtitle = slide.placeholders[1]
    subtitle.text = ("Smart Aquaculture Monitoring Using IoT & Cloud Technology\n\n"
                    "Presented by:\n"
                    "Thirumalasetti Bhogeswara Rao\n"
                    "B.Tech – [Your Branch]\n"
                    "Date: 29-03-2026")
    subtitle.text_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    # 2. Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, WHITE)
    title = slide.shapes.title
    title.text = "Introduction"
    content = slide.placeholders[1]
    content.text = ("• Aquaculture plays a major role in India’s economy.\n"
                    "• Water quality directly impacts fish growth and survival.\n\n"
                    "Traditional fish farming lacks:\n"
                    "• Real-time monitoring\n"
                    "• Automated alerts\n"
                    "• Data-based decision support\n\n"
                    "Aquasphere provides a smart, digital monitoring system for fish ponds.")

    # 3. Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, WHITE)
    title = slide.shapes.title
    title.text = "Problem Statement"
    content = slide.placeholders[1]
    content.text = ("Fish farmers face:\n"
                    "• Manual water testing\n"
                    "• Delayed detection of poor water quality\n"
                    "• Sudden fish mortality\n"
                    "• Lack of technical monitoring tools\n"
                    "• Financial losses\n\n"
                    "There is a need for a low-cost, automated solution.")

    # 4. Proposed Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, LIGHT_BLUE)
    title = slide.shapes.title
    title.text = "Proposed Solution – Aquasphere"
    content = slide.placeholders[1]
    content.text = ("Aquasphere is a:\n"
                    "✔ IoT-based monitoring system\n"
                    "✔ Cloud-connected web application\n"
                    "✔ Real-time data analytics platform\n"
                    "✔ Alert-based notification system\n\n"
                    "It enables farmers to monitor pond conditions remotely.")

    # 5. Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, WHITE)
    title = slide.shapes.title
    title.text = "Objectives"
    content = slide.placeholders[1]
    content.text = ("• Develop a real-time water monitoring system\n"
                    "• Provide web-based dashboard\n"
                    "• Store historical water data\n"
                    "• Generate alerts for abnormal conditions\n"
                    "• Improve fish survival rate")

    # 6. System Architecture Diagram (Text Representation)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, DEEP_SEA)
    title = slide.shapes.title
    title.text = "System Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    content = slide.placeholders[1]
    content.text = ("Water Sensors (pH, Temp, Turbidity)\n"
                    "             ↓\n"
                    "     Microcontroller (IoT Device)\n"
                    "             ↓\n"
                    "         Internet/WiFi\n"
                    "             ↓\n"
                    "        Cloud Server (Render)\n"
                    "             ↓\n"
                    "         Backend API\n"
                    "             ↓\n"
                    "          Database\n"
                    "             ↓\n"
                    "       Web Dashboard (User)")
    content.text_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE

    # 7. Working Process
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, WHITE)
    title = slide.shapes.title
    title.text = "Working Process"
    content = slide.placeholders[1]
    content.text = ("1. Sensors collect water parameters\n"
                    "2. Microcontroller processes data\n"
                    "3. Data sent to cloud server\n"
                    "4. Stored in database\n"
                    "5. Displayed on dashboard\n"
                    "6. Alerts triggered if values exceed threshold")

    # 8. Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, LIGHT_BLUE)
    title = slide.shapes.title
    title.text = "Technology Stack"
    content = slide.placeholders[1]
    content.text = ("Frontend:\n"
                    "• HTML, CSS, JavaScript\n\n"
                    "Backend:\n"
                    "• Python (Flask)\n\n"
                    "Database:\n"
                    "• PostgreSQL (Supabase)\n\n"
                    "Cloud Hosting:\n"
                    "• Render\n\n"
                    "IoT Integration:\n"
                    "• Sensor-based data input")

    # 9. Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, WHITE)
    title = slide.shapes.title
    title.text = "Features"
    content = slide.placeholders[1]
    content.text = ("✔ Real-time monitoring\n"
                    "✔ Secure login system\n"
                    "✔ Data visualization dashboard\n"
                    "✔ Historical data storage\n"
                    "✔ Scalable cloud deployment\n"
                    "✔ Expandable for AI integration")

    # 10. Innovation & Uniqueness
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, WHITE)
    title = slide.shapes.title
    title.text = "Innovation & Uniqueness"
    content = slide.placeholders[1]
    content.text = ("• Designed for Indian aquaculture\n"
                    "• Cost-effective compared to imported systems\n"
                    "• Cloud-based remote access\n"
                    "• Expandable for AI-based predictions\n"
                    "• Scalable for multiple ponds")

    # 11. Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, WHITE)
    title = slide.shapes.title
    title.text = "Future Enhancements"
    content = slide.placeholders[1]
    content.text = ("• Mobile App Integration\n"
                    "• SMS & WhatsApp Alerts\n"
                    "• AI-based disease prediction\n"
                    "• Multi-language interface\n"
                    "• Government scheme integration")

    # 12. Financial Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, LIGHT_BLUE)
    title = slide.shapes.title
    title.text = "Financial Plan (Academic)"
    content = slide.placeholders[1]
    content.text = ("Estimated Cost:\n"
                    "• IoT Hardware Setup – ₹5,000\n"
                    "• Cloud Hosting – ₹2,000/year\n"
                    "• Maintenance – ₹3,000\n\n"
                    "Revenue Model (Future):\n"
                    "• Subscription – ₹499/month\n"
                    "• Hardware Sales – One-time purchase\n"
                    "• Annual Maintenance Plan")

    # 13. Impact
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, WHITE)
    title = slide.shapes.title
    title.text = "Impact"
    content = slide.placeholders[1]
    content.text = ("Aquasphere can:\n"
                    "✔ Reduce fish mortality\n"
                    "✔ Improve farmer income\n"
                    "✔ Promote sustainable aquaculture\n"
                    "✔ Support Digital India\n"
                    "✔ Encourage AgriTech innovation")

    # 14. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide, DEEP_SEA)
    title = slide.shapes.title
    title.text = "Conclusion"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    
    subtitle = slide.placeholders[1]
    subtitle.text = ("Aquasphere combines IoT, cloud computing, and web technologies to modernize aquaculture monitoring.\n\n"
                    "It is scalable, affordable, and impactful.\n\n"
                    "Thank You")
    subtitle.text_frame.paragraphs[0].font.color.rgb = ACCENT

    prs.save("AquaSphere_Startup_Pitch.pptx")
    print("College Demo PPT Generated Successfully!")

if __name__ == "__main__":
    create_college_demo_deck()
